"""Fill the designed PowerPoint template from a fitted MMM — model numbers and
charts go straight into the template's shapes; no AI in this layer.

The template (bundled at ``reporting/deck/templates/report_template.pptx``; see
:func:`default_template_path`) is a finished 24-slide readout. :func:`build_pptx`
reads its slides by stable label text + geometry and
fills the data-bearing ones: the headline KPI cards, the channel scorecard, the
ROI / next-dollar (marginal-ROI) charts, the decomposition, and the per-channel
deep-dives (each with the breakthrough/optimal/saturation zone chart from
:func:`mmm_framework.reporting.helpers.compute_response_zones`). It uses **80%
ranges** to match the template, fills the existing channel rows/slides up to the
model's channel count, and trims the rest.

Per-slide AI insights and the whole-deck synthesis (PR 3) are injected via the
optional ``insights`` map; everything here is deterministic.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

import numpy as np
from pptx.util import Inches

from . import charts, template as T, textfit
from .charts import _money


def default_template_path() -> Path:
    """Path to the bundled deck template.

    Resolution order: the ``MMM_DECK_TEMPLATE`` env override, then the template
    packaged with the library. The agent/API pass an explicit path; this is the
    default used by tools and tests.
    """
    import os

    env = os.environ.get("MMM_DECK_TEMPLATE")
    if env and Path(env).exists():
        return Path(env)
    return Path(__file__).parent / "templates" / "report_template.pptx"


def _pct(samples: np.ndarray, hdi_prob: float) -> tuple[float, float]:
    lo = float(np.percentile(samples, (1 - hdi_prob) / 2 * 100))
    hi = float(np.percentile(samples, (1 + hdi_prob) / 2 * 100))
    return lo, hi


def _read_action(mean: float, lo: float, hi: float, be: float) -> tuple[str, str]:
    """The template's READ / ACTION vocabulary, from the ROI credible interval vs
    the break-even line."""
    if lo > be:
        return "Confidently profitable", "Scale"
    if hi < be:
        return "Below break-even", "Reduce"
    if (hi - lo) > 0.9 * be and mean > 0.9 * be:
        return "High upside, unproven", "Test"
    return "Near break-even", "Hold"


def _half_life_weeks(model: Any, channel: str) -> float | None:
    # This called `_get_adstock_alpha(model, posterior, channel)` — three
    # positional args into a two-arg signature — so it raised TypeError on every
    # call and the deck's carryover half-life was unconditionally None, for
    # geometric channels too. It also derived the horizon from log(0.5)/log(alpha),
    # which is meaningless for a delayed or Weibull kernel.
    try:
        from ...transforms.carryover import (
            carryover_half_life,
            posterior_carryover_kernels,
        )

        k = posterior_carryover_kernels(model, [channel]).get(channel)
        if k is None or not np.all(np.isfinite(k.kernel)):
            return None
        hl = carryover_half_life(k.kernel)
        return float(np.nanmean(hl)) if np.any(np.isfinite(hl)) else None
    except Exception:
        return None


def _portfolio_metrics(
    model: Any, roi_records: list[dict], total_revenue: float | None, hdi_prob: float
) -> dict[str, Any]:
    """Total marketing-attributed revenue, share of revenue, and blended return
    per $1 — point + 80% range, from posterior portfolio-contribution draws."""
    total_spend = float(sum(r.get("spend", 0.0) for r in roi_records))
    portfolio = None
    try:
        cc = model.sample_channel_contributions(max_draws=300)  # (draws, obs, channel)
        portfolio = np.asarray(cc).sum(axis=(1, 2))  # per-draw total media contribution
    except Exception:
        portfolio = None

    out: dict[str, Any] = {"total_spend": total_spend}
    if portfolio is not None and portfolio.size:
        rev_mean = float(portfolio.mean())
        rev_lo, rev_hi = _pct(portfolio, hdi_prob)
        out["revenue"] = (rev_mean, rev_lo, rev_hi)
        if total_spend > 0:
            roi = portfolio / total_spend
            out["blended_roi"] = (float(roi.mean()), *_pct(roi, hdi_prob))
        if total_revenue and total_revenue > 0:
            sh = portfolio / total_revenue
            out["share"] = (float(sh.mean()), *_pct(sh, hdi_prob))
    else:  # point-only fallback
        rev_mean = float(sum(r.get("contribution_mean", 0.0) for r in roi_records))
        out["revenue"] = (rev_mean, float("nan"), float("nan"))
        if total_spend > 0:
            out["blended_roi"] = (rev_mean / total_spend, float("nan"), float("nan"))
        if total_revenue:
            out["share"] = (rev_mean / total_revenue, float("nan"), float("nan"))
    return out


def _standfirst_text(r: dict, z: Any, be: float, currency: str) -> str:
    """Deterministic per-channel standfirst, used when no AI narrative is
    supplied — never leaves the template's example-model prose on a real deck."""
    mroi = z.current_mroi if z is not None else float("nan")
    mroi_s = f"{mroi:.2f}" if np.isfinite(mroi) else "—"
    clause = {
        "Scale": "The case to add budget holds across the plausible range.",
        "Test": "High potential but a wide range — prove it with a controlled test "
        "before scaling.",
        "Hold": "No clear case to scale or cut on current evidence.",
        "Reduce": "On current evidence it returns less than it costs.",
    }[r["action"]]
    return (
        f"{r['channel']} returns {currency}{r['roi_mean']:.2f} per {currency}1 "
        f"(80% range {r['roi_lo']:.2f}–{r['roi_hi']:.2f}); the next dollar returns "
        f"{mroi_s}. {clause}"
    )


def _what_to_do_text(r: dict, z: Any, be: float) -> str:
    """Deterministic "What to do" advice grounded in the channel's own numbers
    (action, marginal return, carryover) — replaces the template example."""
    hl = r.get("half_life")
    mroi = z.current_mroi if z is not None else float("nan")
    action = r["action"]
    if action == "Scale":
        # Scale is assigned from the AVERAGE-ROI interval; the marginal dollar
        # may still be under water — never claim it clears break-even unless it does
        if np.isfinite(mroi) and mroi >= be:
            txt = "Increase weight stepwise and re-check saturation after each step — the marginal dollar still clears break-even."
        else:
            mroi_s = f" (currently {mroi:.2f})" if np.isfinite(mroi) else ""
            txt = (
                "Scale carefully: the average return is confidently profitable, but "
                f"the marginal dollar is already below break-even{mroi_s} — favour "
                "efficiency moves over straight budget increases."
            )
    elif action == "Test":
        txt = "Fund at the current level and run a matched-market test — the range is too wide to commit more on model evidence alone."
    elif action == "Hold":
        if hl is not None and hl > 4:
            txt = (
                f"Hold spend and keep it continuous — the ~{hl:.0f}-week carryover "
                "means gaps waste effect already paid for."
            )
        else:
            txt = "Hold spend at the current level and revisit when a test or new data moves the estimate — decay is quick enough to re-flight freely."
    else:  # Reduce
        mroi_s = f"{mroi:.2f}" if np.isfinite(mroi) else "below break-even"
        txt = (
            f"Step spend down toward the efficient range and redeploy the savings — "
            f"the marginal dollar returns {mroi_s}."
        )
    return txt


def _headline_text(pf: dict, rows: list[dict], currency: str) -> str | None:
    """Deterministic S2 headline when no AI synthesis is supplied."""
    if "share" in pf and np.isfinite(pf["share"][0]):
        share = pf["share"][0]
        if "blended_roi" in pf and np.isfinite(pf["blended_roi"][0]):
            return (
                f"Marketing drives {share:.0%} of revenue at "
                f"{currency}{pf['blended_roi'][0]:.2f} back per {currency}1"
            )
        return f"Marketing drives {share:.0%} of revenue"
    if "blended_roi" in pf and np.isfinite(pf["blended_roi"][0]):
        return f"Marketing returns {currency}{pf['blended_roi'][0]:.2f} per {currency}1"
    return None


def _headline_standfirst(pf: dict, rows: list[dict], be: float) -> str | None:
    """Deterministic S2 standfirst: where the blended return sits and the top
    reallocation move implied by the scorecard.

    "Move budget toward X" is only claimed when X actually earns it (Scale, or
    a profitable Test candidate) — never toward the least-bad channel of an
    all-below-break-even portfolio, which would contradict the scorecard."""
    if not rows:
        return None
    worst = min(rows, key=lambda r: r["roi_mean"])
    parts = []
    if "blended_roi" in pf and np.isfinite(pf["blended_roi"][0]):
        m = pf["blended_roi"][0]
        rel = "above" if m > be * 1.02 else "below" if m < be * 0.98 else "at"
        parts.append(f"The blended return sits {rel} break-even.")
    fundable = [r for r in rows if r["action"] == "Scale"] or [
        r for r in rows if r["action"] == "Test" and r["roi_mean"] > be
    ]
    if fundable:
        best = max(fundable, key=lambda r: r["roi_mean"])
        if best is not worst:
            toward = (
                f"toward {best['channel']} ({best['roi_mean']:.2f})"
                if best["action"] == "Scale"
                else f"toward testing {best['channel']} ({best['roi_mean']:.2f})"
            )
            parts.append(
                f"The gain is not in spending more — it is in moving budget from "
                f"{worst['channel']} ({worst['roi_mean']:.2f}) {toward}."
            )
    else:
        parts.append(
            f"No channel clears break-even with confidence on current evidence — "
            f"the first move is trimming {worst['channel']} "
            f"({worst['roi_mean']:.2f}) and funding tests before adding budget."
        )
    return " ".join(parts) or None


def _cluster_rows(shapes, gap_in: float = 0.3) -> list[list]:
    """Group shapes into visual rows by their top coordinate."""
    from pptx.util import Inches

    if not shapes:
        return []
    shapes = sorted(shapes, key=lambda s: T._emu(s.top))
    gap = int(Inches(gap_in))
    rows, cur, last = [], [shapes[0]], T._emu(shapes[0].top)
    for sh in shapes[1:]:
        t = T._emu(sh.top)
        if t - last > gap:
            rows.append(cur)
            cur = []
        cur.append(sh)
        last = t
    rows.append(cur)
    return rows


def _fill_scorecard(
    slide, rows: list[dict], currency: str, be: float, styles: dict | None = None
) -> list:
    """Fill the channel scorecard (positional: columns by header left, rows by
    top), filling one model channel per template row and blanking the rest.

    Recolours each row's READ pill + ACTION text to the *filled* status's
    design colours (the template rows are coloured for the example model).
    Returns the READ pill text shapes so the pill backgrounds can be resized
    after the global text-fit pass."""
    from pptx.util import Inches

    cols = {}
    for key, label in (
        ("channel", "CHANNEL"),
        ("spend", "SPEND"),
        ("return", "RETURN / $1"),
        ("range", "80% RANGE"),
        ("read", "READ"),
        ("action", "ACTION"),
    ):
        sh = T.find_by_label(slide, label)
        if sh is not None:
            cols[key] = T._emu(sh.left)
    head = T.find_by_label(slide, "CHANNEL")
    if head is None or "channel" not in cols:
        return []
    header_top = T._emu(head.top)
    tol = int(Inches(0.6))

    # data text shapes below the header, excluding the full-width footer note
    data = [
        sh
        for sh in T.iter_text_shapes(slide)
        if T._emu(sh.top) > header_top + int(Inches(0.3))
        and T._emu(sh.width) < int(Inches(12))
    ]
    row_groups = _cluster_rows(data)

    pill_shapes = []
    for i, group in enumerate(row_groups):
        # assign each shape in the row to its nearest column
        cell = {}
        for sh in group:
            left = T._emu(sh.left)
            best = min(cols, key=lambda k: abs(cols[k] - left))
            if abs(cols[best] - left) <= tol and best not in cell:
                cell[best] = sh
        if i < len(rows):
            r = rows[i]
            mean, lo, hi = r["roi_mean"], r["roi_lo"], r["roi_hi"]
            read, action = _read_action(mean, lo, hi, be)
            if "channel" in cell:
                T.set_text(cell["channel"], r["channel"])
            if "spend" in cell:
                T.set_text(cell["spend"], _money(r["spend"], currency))
            if "return" in cell:
                T.set_text(cell["return"], f"{mean:.2f}")
            if "range" in cell:
                T.set_text(cell["range"], f"{lo:.2f} – {hi:.2f}")
            if "read" in cell:
                T.set_text(cell["read"], read)
                if styles:
                    T.restyle_status_shape(slide, cell["read"], read, styles)
                pill_shapes.append(cell["read"])
            if "action" in cell:
                T.set_text(cell["action"], action)
                if styles:
                    T.restyle_status_shape(slide, cell["action"], action, styles)
        else:  # blank the template's extra rows AND its now-orphaned pill decor
            for sh in group:
                if T._norm(sh.text_frame.text) in T.READ_TO_ACTION:
                    bg, dot = T.find_pill_parts(slide, sh)
                    for decor in (bg, dot):
                        if decor is not None:
                            T.delete_shape(decor)
                T.set_text(sh, "")
    return pill_shapes


def _fill_channel_slide(
    slide,
    r: dict,
    z: Any,
    currency: str,
    be: float,
    narrative: str | None,
    styles: dict | None = None,
) -> list:
    """Fill one per-channel deep-dive slide (channel name, action pill, the five
    metric cards, the narrative standfirst, the "What to do" advice, and the
    saturation/zone chart). Returns pill text shapes for post-fit resizing."""
    pill_shapes = []
    # channel name = the prominent text near the top-left
    for sh in T.iter_text_shapes(slide):
        if T._emu(sh.top) < int(Inches(1.2)) and T._emu(sh.left) < int(Inches(2.0)):
            T.set_text(sh, r["channel"])
            break

    _, action = _read_action(r["roi_mean"], r["roi_lo"], r["roi_hi"], be)
    for sh in T.iter_text_shapes(slide):  # action pill (Scale/Test/Hold/Reduce)
        if T._norm(sh.text_frame.text) in ("scale", "test", "hold", "reduce"):
            T.set_text(sh, action)
            if styles:
                T.restyle_status_shape(slide, sh, action, styles)
            pill_shapes.append(sh)
            break

    mroi = z.current_mroi if z is not None else float("nan")
    T.fill_card(
        slide,
        "RETURN / $1",
        f"{r['roi_mean']:.2f}",
        f"80% {r['roi_lo']:.2f}–{r['roi_hi']:.2f}",
    )
    T.fill_card(slide, "CONTRIBUTION", _money(r["contribution"], currency))
    T.fill_card(slide, "SPEND", _money(r["spend"], currency))
    T.fill_card(
        slide,
        "MARGINAL / $1",
        "—" if not np.isfinite(mroi) else f"{mroi:.2f}",
        (
            "clears break-even"
            if (np.isfinite(mroi) and mroi >= be)
            else "below break-even"
        ),
    )
    hl = r.get("half_life")
    T.fill_card(
        slide,
        "CARRYOVER HALF-LIFE",
        "—" if hl is None else f"{hl:.1f}w",
        (
            "fast decay"
            if (hl or 0) < 1.5
            else "slow decay" if (hl or 0) > 4 else "medium decay"
        ),
    )

    # the wide standfirst near the top: AI narrative when supplied, else a
    # deterministic sentence — never the template's example-model prose
    standfirst = narrative or _standfirst_text(r, z, be, currency)
    for sh in T.iter_text_shapes(slide):
        if T._emu(sh.top) < int(Inches(2.2)) and T._emu(sh.width) > int(Inches(8)):
            T.set_text(sh, standfirst)
            break

    # the "What to do" advice line at the bottom — grounded in this channel's
    # numbers (the template's example advice can contradict the filled data)
    wtd = T.find_by_prefix(slide, "What to do")
    if wtd is not None:
        T.set_body_after_label(wtd, _what_to_do_text(r, z, be))

    # the saturation/zone chart on the left panel — rendered to fit the box's
    # exact aspect ratio (no squish).
    if z is not None:
        try:
            T.replace_image_fit(
                slide,
                lambda w, h: charts.saturation_zones_png(
                    z, currency=currency, width=w, height=h
                ),
                match=T.pictures_in_region(slide, Inches(1.05), Inches(4.43), 0, 0),
            )
        except Exception:
            pass
    return pill_shapes


def _recommend_moves(rows: list[dict], currency: str) -> list[tuple[str, str]]:
    """Up to three (lead, body) recommendation cards for the S2 "WHAT WE
    RECOMMEND" band, from the scorecard's action buckets."""

    def _names(chs):
        return " · ".join(c["channel"] for c in chs[:3])

    def _fmt(chs):
        return ", ".join(f"{c['channel']} {c['roi_mean']:.2f}" for c in chs[:3])

    buckets = {
        a: sorted((r for r in rows if r["action"] == a), key=lambda r: -r["roi_mean"])
        for a in ("Scale", "Reduce", "Test", "Hold")
    }
    cards: list[tuple[str, str]] = []
    if buckets["Scale"]:
        chs = buckets["Scale"]
        cards.append(
            (
                f"Scale {_names(chs)} ",
                f"— the whole credible range clears break-even ({_fmt(chs)}); "
                "additional budget keeps earning.",
            )
        )
    if buckets["Reduce"]:
        chs = buckets["Reduce"]
        spend = sum(c["spend"] for c in chs)
        cards.append(
            (
                f"Reduce {_names(chs)} ",
                f"— {_money(spend, currency)} of spend returning under a dollar "
                f"across the range ({_fmt(chs)}); trim and redeploy.",
            )
        )
    if buckets["Test"]:
        chs = buckets["Test"]
        cards.append(
            (
                f"Test {_names(chs)} ",
                f"— attractive central returns ({_fmt(chs)}) but ranges too wide "
                "to fund on faith; prove them with controlled tests.",
            )
        )
    if len(cards) < 3 and buckets["Hold"]:
        chs = buckets["Hold"]
        cards.append(
            (
                f"Hold {_names(chs)} ",
                f"— near break-even ({_fmt(chs)}); keep steady until a test or "
                "more data moves the estimate.",
            )
        )
    return cards[:3]


def _fill_recommend_cards(slide, rows: list[dict], currency: str) -> None:
    """Fill the S2 "WHAT WE RECOMMEND" band's three cards from the model —
    the template ships example-model recommendations ("Scale Video …") that
    would otherwise survive into every client deck."""
    anchor = T.find_by_label(slide, "WHAT WE RECOMMEND")
    if anchor is None:
        return
    a_bot = T._emu(anchor.top) + T._emu(anchor.height)
    cards = [
        sh
        for sh in T.iter_text_shapes(slide)
        if T._emu(sh.top) > a_bot and T._emu(sh.top) < a_bot + int(Inches(1.0))
    ]
    cards.sort(key=lambda sh: T._emu(sh.left))
    recs = _recommend_moves(rows, currency)
    for i, sh in enumerate(cards[:3]):
        if i < len(recs):
            lead, body = recs[i]
            if not T.set_body_after_label(sh, body):
                T.set_text(sh, lead + body)
                continue
            sh.text_frame.paragraphs[0].runs[0].text = lead
        else:
            T.delete_backing_card(slide, sh)
            T.set_text(sh, "")


def _fill_moves_slide(slide, rows: list[dict], currency: str, be: float) -> None:
    """Fill the "reallocation in four moves" quadrants (S10) with the model's
    actual channel buckets — the template ships example-model channels."""

    def _fmt(chs: list[dict]) -> str:
        return ", ".join(f"{c['channel']} {c['roi_mean']:.2f}" for c in chs)

    buckets = {
        a: [r for r in rows if r["action"] == a]
        for a in ("Scale", "Test", "Hold", "Reduce")
    }
    rationale = {
        "Scale": lambda chs: (
            f"The credible range clears break-even ({_fmt(chs)}) and the marginal "
            "dollar still earns — additional budget keeps paying back."
        ),
        "Test": lambda chs: (
            f"Attractive central returns ({_fmt(chs)}) but the ranges are too wide "
            "to fund on faith — a controlled test settles it."
        ),
        "Hold": lambda chs: (
            f"Near break-even ({_fmt(chs)}) with no clear case to scale or cut — "
            "keep steady until a test or more data moves the estimate."
        ),
        "Reduce": lambda chs: (
            f"Below break-even across the plausible range ({_fmt(chs)}) — trim "
            "toward the efficient level and redeploy the savings."
        ),
    }
    for action, chs in buckets.items():
        anchor = T.find_by_label(slide, action)
        if anchor is None:
            continue
        below = T.shapes_below(slide, anchor, left_tol_in=0.4, max_n=2)
        if len(below) < 2:
            continue
        if chs:
            T.set_text(below[0], " · ".join(c["channel"] for c in chs))
            T.set_text(below[1], rationale[action](chs))
        else:
            T.set_text(below[0], "—")
            T.set_text(below[1], "No channels in this bucket this cycle.")


def _recommend_tests(rows: list[dict]) -> list[tuple[str, str]]:
    """Up to three deterministic experiment recommendations from the scorecard:
    prove the highest unproven upside, de-risk the biggest cut, then tighten
    the next-widest range."""
    recs: list[tuple[str, str]] = []
    used: set[str] = set()
    tests = sorted(
        (r for r in rows if r["action"] == "Test"), key=lambda r: -r["roi_hi"]
    )
    reduces = sorted(
        (r for r in rows if r["action"] == "Reduce"), key=lambda r: -r["spend"]
    )
    if tests:
        r = tests[0]
        used.add(r["channel"])
        recs.append(
            (
                f"Geo holdout on {r['channel']}",
                f"Highest unproven upside (central {r['roi_mean']:.2f}, range "
                f"{r['roi_lo']:.2f}–{r['roi_hi']:.2f}). A matched-market holdout "
                "settles whether to scale or stand down.",
            )
        )
    if reduces:
        r = reduces[0]
        used.add(r["channel"])
        recs.append(
            (
                f"Spend-down test on {r['channel']}",
                f"Confirms {r['channel']} can be cut without losing revenue before "
                f"budget moves — its range sits below break-even (central "
                f"{r['roi_mean']:.2f}).",
            )
        )
    # widest remaining credible interval = the read most worth buying
    rest = sorted(
        (r for r in rows if r["channel"] not in used),
        key=lambda r: -(r["roi_hi"] - r["roi_lo"]),
    )
    for r in rest:
        if len(recs) >= 3:
            break
        recs.append(
            (
                f"Incrementality test on {r['channel']}",
                f"Central {r['roi_mean']:.2f} with an 80% range of "
                f"{r['roi_lo']:.2f}–{r['roi_hi']:.2f}; a clean read decides "
                f"whether {r['channel']} earns a bigger role in the next plan.",
            )
        )
    return recs[:3]


def _fill_tests_slide(slide, rows: list[dict]) -> None:
    """Fill the numbered recommended-tests triplets (S22) from the model."""
    recs = _recommend_tests(rows)
    for i in range(3):
        num = T.find_by_label(slide, f"{i + 1:02d}")
        if num is None:
            continue
        n_top, n_right = T._emu(num.top), T._emu(num.left) + T._emu(num.width)
        title = body = None
        for sh in T.iter_text_shapes(slide):
            if sh is num or T._emu(sh.left) < n_right:
                continue
            if abs(T._emu(sh.top) - n_top) <= int(Inches(0.15)):
                title = sh
        if title is not None:
            below = T.shapes_below(slide, title, left_tol_in=0.3, max_n=1)
            body = below[0] if below else None
        if i < len(recs):
            t, b = recs[i]
            if title is not None:
                T.set_text(title, t)
            if body is not None:
                T.set_text(body, b)
        else:
            # blank the unused triplet AND its card frame (else an empty
            # outlined card is left behind)
            T.delete_backing_card(slide, num)
            for sh in (num, title, body):
                if sh is not None:
                    T.set_text(sh, "")


def build_pptx(
    model: Any,
    *,
    template_path: str | Path | None = None,
    out_path: str | Path | None = None,
    deck: Any = None,
    client: str | None = None,
    kpi_name: str = "Revenue",
    currency: str = "$",
    break_even: float = 1.0,
    margin: float | None = None,
    hdi_prob: float = 0.8,
    max_channels: int = 7,
    insights: dict[str, str] | None = None,
) -> bytes:
    """Fill the template deck from a fitted model and return the .pptx bytes
    (also written to ``out_path`` if given).

    ``hdi_prob`` defaults to 0.8 to match the template's "80% range". ``margin``
    sets a profit-maximizing break-even (1/margin). ``insights`` (PR 3) maps slide
    keys to AI narrative text; omitted here.
    """
    from pptx import Presentation

    from ..helpers import compute_response_zones, compute_roi_with_uncertainty

    template_path = template_path or default_template_path()
    eff_be = (1.0 / float(margin)) if margin else float(break_even)

    roi_df = compute_roi_with_uncertainty(model, hdi_prob=hdi_prob)
    roi_records = (
        roi_df.to_dict("records") if roi_df is not None and len(roi_df) else []
    )
    zones = {}
    try:
        zones = compute_response_zones(model, break_even=eff_be, hdi_prob=hdi_prob)
    except Exception:
        zones = {}

    bundle = None
    total_revenue = None
    try:
        from ..extractors import create_extractor

        bundle = create_extractor(model).extract()
        if getattr(bundle, "actual", None) is not None:
            total_revenue = float(np.asarray(bundle.actual).sum())
    except Exception:
        bundle = None

    # per-channel rows (sorted by action priority then ROI)
    _ORDER = {"Scale": 0, "Test": 1, "Hold": 2, "Reduce": 3}
    rows = []
    for r in roi_records:
        ch = r["channel"]
        lo = r.get("roi_hdi_low", r.get("roi_mean", 0.0))
        hi = r.get("roi_hdi_high", r.get("roi_mean", 0.0))
        _, action = _read_action(r.get("roi_mean", 0.0), lo, hi, eff_be)
        rows.append(
            {
                "channel": ch,
                "spend": r.get("spend", 0.0),
                "contribution": r.get("contribution_mean", 0.0),
                "roi_mean": r.get("roi_mean", 0.0),
                "roi_lo": lo,
                "roi_hi": hi,
                "action": action,
                "half_life": _half_life_weeks(model, ch),
            }
        )
    rows.sort(key=lambda r: (_ORDER.get(r["action"], 9), -r["roi_mean"]))
    rows = rows[:max_channels]

    pf = _portfolio_metrics(model, roi_records, total_revenue, hdi_prob)
    insights = insights or {}

    prs = Presentation(str(template_path))
    slides = list(prs.slides)
    # per-status design colours, read from the template's own pills before any
    # fill overwrites them
    styles = T.harvest_status_styles(prs)
    pills: list[tuple[Any, Any]] = []  # (slide, pill text shape) to resize later

    # ---- S0: cover ----
    if len(slides) > 0:
        s = slides[0]
        sub = T.find_by_prefix(s, "Prepared for")
        if sub is not None and client:
            T.set_text(sub, f"Prepared for {client} · Planning & analytics")

    # ---- S1: headline KPI cards ----
    if len(slides) > 1:
        s = slides[1]
        if "revenue" in pf:
            m, lo, hi = pf["revenue"]
            rng = (
                ""
                if not np.isfinite(lo)
                else f"80% range {_money(lo, currency)} – {_money(hi, currency)}"
            )
            T.fill_card(s, "MARKETING-ATTRIBUTED REVENUE", _money(m, currency), rng)
        if "share" in pf:
            m, lo, hi = pf["share"]
            rng = "" if not np.isfinite(lo) else f"80% range {lo:.1%} – {hi:.1%}"
            T.fill_card(s, "SHARE OF TOTAL REVENUE", f"{m:.1%}", rng)
        if "blended_roi" in pf:
            m, lo, hi = pf[
                "blended_roi"
            ]  # a ratio ($ returned per $1) — 2 decimals, not abbreviated
            rng = (
                ""
                if not np.isfinite(lo)
                else f"80% range {currency}{lo:.2f} – {currency}{hi:.2f}"
            )
            T.fill_card(s, "BLENDED RETURN PER $1", f"{currency}{m:.2f}", rng)
        # headline + standfirst: AI insights when supplied, else a deterministic
        # summary — the template's example-model claims never survive a fill
        headline = insights.get("headline") or _headline_text(pf, rows, currency)
        standfirst = insights.get("standfirst") or _headline_standfirst(
            pf, rows, eff_be
        )
        h = T.find_by_label(s, "THE HEADLINE")
        if h is not None:
            # below the eyebrow: [0] = the big title, [1] = the standfirst para
            below = T.shapes_below(s, h, left_tol_in=2.0, max_n=2)
            if below and headline:
                T.set_text(below[0], headline)
            if len(below) > 1 and standfirst:
                T.set_text(below[1], standfirst)
        if rows:
            _fill_recommend_cards(s, rows, currency)

    # ---- S5: channel scorecard ----
    sc = next((s for s in slides if T.find_by_label(s, "CHANNEL SCORECARD")), None)
    if sc is not None:
        pills += [
            (sc, sh) for sh in _fill_scorecard(sc, rows, currency, eff_be, styles)
        ]

    # ---- S6 / S7: ROI and marginal-ROI forests ----
    roi_dict = {
        r["channel"]: {
            "mean": r["roi_mean"],
            "lower": r["roi_lo"],
            "upper": r["roi_hi"],
        }
        for r in rows
    }
    s6 = next((s for s in slides if T.find_by_label(s, "RETURN & UNCERTAINTY")), None)
    if s6 is not None and roi_dict:
        try:
            T.replace_image_fit(
                s6,
                lambda w, h: charts.roi_forest_png(
                    roi_dict, break_even=eff_be, width=w, height=h
                ),
            )
        except Exception:
            pass
    s7 = next((s for s in slides if T.find_by_label(s, "THE NEXT DOLLAR")), None)
    if s7 is not None and zones:
        mroi_dict = {
            ch: {
                "mean": z.current_mroi,
                "lower": z.current_mroi_hdi[0],
                "upper": z.current_mroi_hdi[1],
            }
            for ch, z in zones.items()
        }
        try:
            T.replace_image_fit(
                s7,
                lambda w, h: charts.roi_forest_png(
                    mroi_dict,
                    break_even=eff_be,
                    width=w,
                    height=h,
                    title="Marginal return on the next dollar",
                    xlabel="Marginal ROI (next-dollar return)",
                ),
            )
        except Exception:
            pass

    # ---- S4: decomposition (waterfall; hatched control/seasonal swings) ----
    if bundle is not None and getattr(bundle, "component_totals", None):
        s4 = next(
            (
                s
                for s in slides
                if T.find_by_prefix(s, "The full revenue decomposition")
            ),
            None,
        )
        if s4 is not None:
            comp = bundle.component_totals
            comp_series = getattr(bundle, "component_time_series", None)
            media_keys = list(getattr(bundle, "channel_names", None) or []) or [
                r["channel"] for r in roi_records
            ]
            try:
                T.replace_image_fit(
                    s4,
                    lambda w, h: charts.decomposition_png(
                        comp,
                        series=comp_series,
                        media_keys=media_keys,
                        currency=currency,
                        total_label=f"Total {kpi_name.lower()}",
                        kpi_name=kpi_name,
                        width=w,
                        height=h,
                    ),
                )
            except Exception:
                pass

    # ---- S10 / S22: reallocation moves + recommended tests ----
    s10 = next((s for s in slides if T.find_by_label(s, "WHERE TO MOVE BUDGET")), None)
    if s10 is not None and rows:
        _fill_moves_slide(s10, rows, currency, eff_be)
    s22 = next((s for s in slides if T.find_by_label(s, "RECOMMENDED TESTS")), None)
    if s22 is not None and rows:
        _fill_tests_slide(s22, rows)

    # ---- S11: flighting plan — neutralize the example-model prose (the chart
    # itself stays illustrative; the standfirst must not name example channels)
    s11 = next((s for s in slides if T.find_by_label(s, "THE FLIGHTING PLAN")), None)
    if s11 is not None:
        sf = T.find_by_prefix(s11, "The recommended plan")
        if sf is not None:
            T.set_text(
                sf,
                "The recommended plan smooths burst-heavy flighting into steadier "
                "weekly weight and protects continuity on carryover channels — at "
                "the same annual budget. Weekly shapes are illustrative.",
            )

    # ---- S12-18: per-channel deep-dive slides (fill N, delete extras) ----
    deep = [
        (idx, s)
        for idx, s in enumerate(slides)
        if T.find_by_label(s, "RETURN / $1") is not None
        and T.find_by_label(s, "CARRYOVER HALF-LIFE") is not None
    ]
    used_idx = set()
    for j, (idx, s) in enumerate(deep):
        if j < len(rows):
            r = rows[j]
            pills += [
                (s, sh)
                for sh in _fill_channel_slide(
                    s,
                    r,
                    zones.get(r["channel"]),
                    currency,
                    eff_be,
                    insights.get(f"channel:{r['channel']}"),
                    styles,
                )
            ]
            used_idx.add(idx)
    # delete the unused deep-dive slides (highest index first to keep indices valid)
    for idx, _ in sorted(deep, key=lambda t: -t[0]):
        if idx not in used_idx:
            T.delete_slide(prs, idx)

    # ---- final layout passes (order matters) ----
    # 1. repair run-together label runs (two template paragraphs were authored
    #    without their separating space);
    # 2. make every text shape hold its text explicitly — widen single-line
    #    labels / set reduced run sizes — because Google Slides and friends
    #    never recompute PPT's shrink-on-overflow autofit;
    # 3. resize pill backgrounds to hug their (possibly grown) status text;
    # 4. drop the cached autofit scale so PowerPoint re-fits with actual fonts.
    for s in prs.slides:
        T.normalize_run_boundaries(s)
    textfit.fit_presentation_text(prs)
    for s, sh in pills:
        T.hug_pill_background(s, sh)
    for s in prs.slides:
        T.clear_autofit_scale(s)

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    if out_path is not None:
        Path(out_path).write_bytes(data)
    return data


__all__ = ["build_pptx", "default_template_path"]
