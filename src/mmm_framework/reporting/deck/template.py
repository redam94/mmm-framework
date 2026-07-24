"""Low-level python-pptx helpers for filling a *designed* template deck.

The template is a finished, hand-laid-out deck (named ``Text``/``Shape``/``Image``
auto-shapes with example content), not a placeholder layout. So filling it means:
locate a shape by its **stable label text** (e.g. ``"RETURN / $1"``) or by
geometry relative to a label, replace the text **preserving the template's run
formatting**, swap example chart images for model-rendered PNGs, and trim the
template's extra channel rows/slides down to the model's channel count.

python-pptx is imported lazily so importing :mod:`reporting.deck` never requires
it — only :func:`mmm_framework.reporting.deck.builder.build_pptx` does.
"""

from __future__ import annotations

import io
from typing import Any, Callable


def _norm(s: str) -> str:
    return " ".join((s or "").split()).strip().lower()


def shape_text(shape) -> str:
    """The shape's text, or '' if it has none."""
    if shape.has_text_frame:
        return shape.text_frame.text or ""
    return ""


def set_text(shape, text: str) -> bool:
    """Replace a shape's text while keeping the template's formatting.

    Keeps the first paragraph's first run (its font / size / color / bold) and
    rewrites only its characters, deleting any other runs and paragraphs — so the
    designed look is preserved. Returns ``False`` if the shape can't hold text.
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    paras = tf.paragraphs
    if not paras:
        return False
    first = paras[0]
    # delete extra paragraphs (keep paragraph[0]'s XML element)
    for extra in list(paras[1:]):
        extra._p.getparent().remove(extra._p)
    runs = first.runs
    if runs:
        runs[0].text = str(text)
        for r in list(runs[1:]):
            r._r.getparent().remove(r._r)
    else:
        first.add_run().text = str(text)
    return True


def iter_text_shapes(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            yield sh


def find_by_label(slide, label: str):
    """First shape whose (normalized) text equals ``label``."""
    target = _norm(label)
    for sh in iter_text_shapes(slide):
        if _norm(sh.text_frame.text) == target:
            return sh
    return None


def find_by_prefix(slide, prefix: str):
    """First shape whose (normalized) text starts with ``prefix``."""
    target = _norm(prefix)
    for sh in iter_text_shapes(slide):
        if _norm(sh.text_frame.text).startswith(target):
            return sh
    return None


def _emu(v) -> int:
    return int(v) if v is not None else 0


def shapes_below(slide, anchor, *, left_tol_in: float = 0.6, max_n: int = 4) -> list:
    """Text shapes positioned directly below ``anchor`` and roughly left-aligned
    with it, ordered top-to-bottom — i.e. the value/sub lines of a KPI card whose
    label is ``anchor``.
    """
    from pptx.util import Inches

    a_left = _emu(anchor.left)
    a_top = _emu(anchor.top)
    tol = int(Inches(left_tol_in))
    out = []
    for sh in iter_text_shapes(slide):
        if sh is anchor:
            continue
        if abs(_emu(sh.left) - a_left) <= tol and _emu(sh.top) > a_top:
            out.append(sh)
    out.sort(key=lambda s: _emu(s.top))
    return out[:max_n]


def fill_card(slide, label: str, value: str | None, sub: str | None = None) -> bool:
    """Fill a labeled KPI card: the shape with text ``label`` keeps its label;
    the next shape below becomes ``value`` and the one after becomes ``sub``.
    """
    anchor = find_by_label(slide, label)
    if anchor is None:
        return False
    below = shapes_below(slide, anchor)
    if value is not None and len(below) >= 1:
        set_text(below[0], value)
    if sub is not None and len(below) >= 2:
        set_text(below[1], sub)
    return True


def replace_image(
    slide,
    png: bytes,
    *,
    match: Callable[[Any], bool] | None = None,
    region: tuple[float, float, float, float] | None = None,
) -> bool:
    """Replace the picture(s) in a region with ``png``.

    Finds picture shapes (optionally filtered by ``match``), records the geometry
    of the first, removes them all (the template stacks 2–3 layered copies per
    chart), and adds a single new picture at that geometry. ``region`` (EMU
    left/top/width/height) overrides the geometry. Returns ``False`` if no
    picture matched.
    """
    pics = [
        sh
        for sh in slide.shapes
        if sh.shape_type is not None
        and "PICTURE" in str(sh.shape_type)
        and (match is None or match(sh))
    ]
    if not pics:
        return False
    first = pics[0]
    geom = region or (first.left, first.top, first.width, first.height)
    for sh in pics:
        sh._element.getparent().remove(sh._element)
    slide.shapes.add_picture(
        io.BytesIO(png), geom[0], geom[1], width=geom[2], height=geom[3]
    )
    return True


def replace_image_fit(
    slide, render_fn, *, match: Callable[[Any], bool] | None = None
) -> bool:
    """Replace the picture(s) in a region with a chart rendered to **fit the box
    exactly** — no aspect distortion.

    Unlike :func:`replace_image` (which forces an existing PNG into the box's
    width×height and squishes it when the aspects differ), this measures the
    matched box and calls ``render_fn(width_in, height_in)`` to render the chart
    at that exact aspect ratio, then drops it in. Returns ``False`` if no picture
    matched.
    """
    from pptx.util import Emu

    pics = [
        sh
        for sh in slide.shapes
        if sh.shape_type is not None
        and "PICTURE" in str(sh.shape_type)
        and (match is None or match(sh))
    ]
    if not pics:
        return False
    first = pics[0]
    left, top, width, height = first.left, first.top, first.width, first.height
    png = render_fn(Emu(width).inches, Emu(height).inches)
    for sh in pics:
        sh._element.getparent().remove(sh._element)
    slide.shapes.add_picture(io.BytesIO(png), left, top, width=width, height=height)
    return True


def clear_autofit_scale(slide) -> int:
    """Strip the cached shrink-to-fit scale from every text box on the slide.

    The template's text boxes use "shrink text on overflow" with a ``fontScale``
    cached for the *authoring* font. When PowerPoint substitutes a missing font
    (e.g. Calibri), that stale scale is too large and text overflows / wraps a
    character onto the next line. Removing ``fontScale``/``lnSpcReduction`` (while
    keeping autofit on) makes PowerPoint recompute the fit for the actual font on
    open. Returns the number of boxes adjusted.
    """
    from pptx.oxml.ns import qn

    n = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        body = sh.text_frame._txBody.find(qn("a:bodyPr"))
        if body is None:
            continue
        na = body.find(qn("a:normAutofit"))
        if na is None:
            continue
        changed = False
        for attr in ("fontScale", "lnSpcReduction"):
            if attr in na.attrib:
                del na.attrib[attr]
                changed = True
        n += int(changed)
    return n


def pictures_in_region(slide, left, top, width, height, *, tol_in: float = 0.3):
    """Predicate factory: matches pictures whose top-left is within ``tol_in`` of
    a target region's top-left (to target one chart when a slide has several)."""
    from pptx.util import Inches

    tol = int(Inches(tol_in))
    lo_l, lo_t = int(left) - tol, int(top) - tol
    hi_l, hi_t = int(left) + tol, int(top) + tol

    def _match(sh):
        return lo_l <= _emu(sh.left) <= hi_l and lo_t <= _emu(sh.top) <= hi_t

    return _match


def normalize_run_boundaries(slide) -> int:
    """Insert a missing space at run boundaries that would render run-together.

    The designed template styles inline labels as separate runs (e.g. a small
    bold ``"What to do "`` followed by the body text). Two of those label runs
    were authored *without* their trailing space, so every renderer shows
    ``"What to doHold spend…"``. This inserts a single space wherever two
    adjacent runs of **different font size** butt a letter/digit directly
    against a letter/digit — which can only ever be a missing separator, never
    intentional intra-word styling at equal size. Returns the number of fixes.
    """
    n = 0
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        for para in sh.text_frame.paragraphs:
            runs = para.runs
            for a, b in zip(runs, runs[1:]):
                if not a.text or not b.text:
                    continue
                if (
                    a.text[-1].isalnum()
                    and b.text[0].isalnum()
                    and (a.font.size or 0) != (b.font.size or 0)
                ):
                    a.text = a.text + " "
                    n += 1
    return n


def set_body_after_label(shape, body: str) -> bool:
    """Replace the text *after* an inline label run (e.g. ``"What to do "``)
    while keeping the label and both runs' designed formatting. Extra runs are
    dropped. Returns ``False`` when the shape has no label+body structure."""
    if not shape.has_text_frame:
        return False
    paras = shape.text_frame.paragraphs
    if not paras or len(paras[0].runs) < 2:
        return False
    first = paras[0]
    for extra in list(paras[1:]):
        extra._p.getparent().remove(extra._p)
    label = first.runs[0]
    if label.text and not label.text.endswith(" "):
        label.text = label.text + " "
    first.runs[1].text = str(body)
    for r in list(first.runs[2:]):
        r._r.getparent().remove(r._r)
    return True


# ---------------------------------------------------------------------------
# Status styling (READ/ACTION pills)
# ---------------------------------------------------------------------------

# the template's READ ↔ ACTION vocabulary (see builder._read_action)
READ_TO_ACTION = {
    "confidently profitable": "scale",
    "high upside, unproven": "test",
    "near break-even": "hold",
    "below break-even": "reduce",
}


def _fill_rgb(shape):
    """A shape's solid-fill RGB, or None."""
    try:
        from pptx.enum.dml import MSO_FILL_TYPE

        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            return shape.fill.fore_color.rgb
    except Exception:
        pass
    return None


def _set_fill_rgb(shape, rgb) -> None:
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    except Exception:
        pass


def _first_run_color(shape):
    try:
        return shape.text_frame.paragraphs[0].runs[0].font.color.rgb
    except Exception:
        return None


def set_text_color(shape, rgb) -> None:
    """Set every run's font colour in the shape."""
    if rgb is None or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = rgb
            except Exception:
                pass


def find_pill_parts(slide, text_shape) -> tuple[Any, Any]:
    """The rounded-rect background and dot ellipse that decorate ``text_shape``.

    The template draws each status pill as three stacked shapes: a pale
    ``roundRect`` background, a small ``ellipse`` dot, and the text. Both
    decorations are text-free and geometrically contain / sit inside the text
    shape's row. Returns ``(background, dot)`` — either may be ``None``.
    """
    t_l, t_t = _emu(text_shape.left), _emu(text_shape.top)
    t_b = t_t + _emu(text_shape.height)
    t_cy = (t_t + t_b) / 2
    bg = dot = None
    for sh in slide.shapes:
        if sh is text_shape or (
            getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
        ):
            continue
        left, top = _emu(sh.left), _emu(sh.top)
        right, bottom = left + _emu(sh.width), top + _emu(sh.height)
        if not (top <= t_cy <= bottom):
            continue
        try:
            prst_name = sh._element.spPr.prstGeom.attrib.get("prst")
        except Exception:
            prst_name = None
        w_in = _emu(sh.width) / 914400
        h_in = _emu(sh.height) / 914400
        # pill decorations are row-sized — never the big card behind the table
        if prst_name == "roundRect" and h_in < 0.7 and left <= t_l and right >= t_l:
            bg = sh
        elif prst_name == "ellipse" and w_in < 0.2 and left < t_l:
            dot = sh
    return bg, dot


def harvest_status_styles(prs) -> dict[str, dict[str, Any]]:
    """Read the per-status design (text/dot colour + pill background colour +
    background padding) from the template's own scorecard and deep-dive pills,
    keyed by normalized status/action text."""
    styles: dict[str, dict[str, Any]] = {}
    vocab = set(READ_TO_ACTION) | set(READ_TO_ACTION.values())
    for slide in prs.slides:
        for sh in iter_text_shapes(slide):
            key = _norm(sh.text_frame.text)
            if key not in vocab:
                continue
            bg, dot = find_pill_parts(slide, sh)
            found = {
                "text_rgb": _first_run_color(sh),
                "bg_rgb": _fill_rgb(bg) if bg is not None else None,
                "dot_rgb": _fill_rgb(dot) if dot is not None else None,
            }
            # merge: a later, richer sighting (e.g. a deep-dive pill with a
            # background) fills fields an earlier plain-text sighting lacked
            cur = styles.setdefault(key, found)
            for k, v in found.items():
                if cur.get(k) is None and v is not None:
                    cur[k] = v
    # unify READ and ACTION entries: an action inherits its read-status colours
    for read, action in READ_TO_ACTION.items():
        if read in styles and action not in styles:
            styles[action] = styles[read]
        if action in styles and read not in styles:
            styles[read] = styles[action]
    return styles


def restyle_status_shape(slide, text_shape, status: str, styles: dict) -> None:
    """Recolour a filled READ/ACTION cell (text + pill background + dot) to the
    design colours of its *new* status."""
    st = styles.get(_norm(status))
    if not st:
        return
    set_text_color(text_shape, st.get("text_rgb"))
    bg, dot = find_pill_parts(slide, text_shape)
    if bg is not None and st.get("bg_rgb") is not None:
        _set_fill_rgb(bg, st["bg_rgb"])
    if dot is not None and st.get("dot_rgb") is not None:
        _set_fill_rgb(dot, st["dot_rgb"])


def hug_pill_background(slide, text_shape, pad_in: float = 0.10) -> None:
    """Resize a pill's rounded-rect background to hug its (possibly resized)
    text shape, keeping the background's left edge."""
    from pptx.util import Emu, Inches

    bg, _dot = find_pill_parts(slide, text_shape)
    if bg is None:
        return
    new_right = _emu(text_shape.left) + _emu(text_shape.width) + int(Inches(pad_in))
    bg.width = Emu(max(int(Inches(0.3)), new_right - _emu(bg.left)))


def delete_backing_card(slide, shape, max_h_in: float = 2.0) -> int:
    """Delete the text-free card frame(s) (roundRect/rect decor) drawn behind
    ``shape`` — used when blanking a templated card so no empty outline is left.
    Never touches the big page-panel cards (height filter). Returns count."""
    cx = _emu(shape.left) + _emu(shape.width) // 2
    cy = _emu(shape.top) + _emu(shape.height) // 2
    doomed = []
    for sh in slide.shapes:
        if sh is shape or (
            getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
        ):
            continue
        try:
            prst_name = sh._element.spPr.prstGeom.attrib.get("prst")
        except Exception:
            continue
        if prst_name not in ("roundRect", "rect"):
            continue
        if _emu(sh.height) / 914400 > max_h_in:
            continue
        left, top = _emu(sh.left), _emu(sh.top)
        if left <= cx <= left + _emu(sh.width) and top <= cy <= top + _emu(sh.height):
            doomed.append(sh)
    for sh in doomed:
        delete_shape(sh)
    return len(doomed)


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def delete_slide(prs, index: int) -> None:
    """Remove a slide from the presentation by index (drops the relationship and
    the sldId entry; the part is orphaned, which PowerPoint tolerates)."""
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    if 0 <= index < len(ids):
        xml_slides.remove(ids[index])


__all__ = [
    "shape_text",
    "set_text",
    "iter_text_shapes",
    "find_by_label",
    "find_by_prefix",
    "shapes_below",
    "fill_card",
    "replace_image",
    "replace_image_fit",
    "clear_autofit_scale",
    "pictures_in_region",
    "delete_shape",
    "delete_slide",
    "normalize_run_boundaries",
    "set_body_after_label",
    "READ_TO_ACTION",
    "set_text_color",
    "find_pill_parts",
    "harvest_status_styles",
    "restyle_status_shape",
    "hug_pill_background",
]
