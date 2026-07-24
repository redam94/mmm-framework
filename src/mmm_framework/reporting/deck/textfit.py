"""Renderer-agnostic text fitting for the filled slide deck.

The designed template sizes every text box to its *example* text and relies on
PowerPoint's "shrink text on overflow" (``<a:normAutofit/>``) to absorb any
difference. PowerPoint recomputes that shrink on open — but Google Slides (and
some other renderers) import the deck without recomputing it, so any text that
is wider than its box renders at full size and wraps mid-word ("Searc / h") or
clips ("Medi"). The template's boxes are also cut glyph-exact, so even
*untouched* labels can clip under small font-metric drift.

This module makes the deck self-sufficient: after filling, every text shape is
measured with real font metrics and made to fit **explicitly** —

1. single-line labels that are too wide are *widened* (bounded by the slide
   edge and by neighbouring text shapes so columns never collide), and
2. anything still too big gets an explicit, reduced run font size (``sz``),
   which every renderer honours.

Measurement prefers the template's *actual* fonts when they are installed
(matplotlib's font manager finds Arial/Georgia/Courier New on macOS/Windows);
where they are missing it falls back to matplotlib's bundled DejaVu families
with measured width-calibration factors, so estimates stay honest on servers.
"""

from __future__ import annotations

import math
import re
from functools import lru_cache
from typing import Any

EMU_PER_IN = 914400
# multiplied over measured widths: absorbs renderer hinting/substitution drift
SAFETY = 1.04
# single-line boxes may not grow closer than this to a neighbouring text shape
NEIGHBOR_GAP_IN = 0.08
# ... nor closer than this to the slide edge
EDGE_MARGIN_IN = 0.10
# never shrink text below this fraction of its designed size
MIN_SCALE = 0.55
# baseline-to-baseline line height as a fraction of font size (PPT single spacing)
LINE_FACTOR = 1.2
# a box whose inner height is under this many line heights was designed to hold
# ONE line (labels/values); at or above it, the box was designed to wrap.
# Measured on the bundled template: single-line boxes cluster at ≤ 1.10 line
# heights, designed-to-wrap boxes at ≥ 1.6 (see tests).
SINGLE_LINE_RATIO = 1.35
# multi-line boxes only get adjusted when the estimated text height exceeds the
# box by more than this (estimates are deliberately conservative)
MULTI_TOLERANCE = 1.30
# ... and are then shrunk to leave this much estimated slack
MULTI_TARGET = 1.12

_MONO_PAT = re.compile(r"courier|mono|consolas", re.I)
_SERIF_PAT = re.compile(r"georgia|times|serif|garamond|cambria|playfair|didot", re.I)

# width calibration for the bundled DejaVu stand-ins when the template's real
# font is not installed (measured real-font ÷ DejaVu ratios; see tests)
_DEJAVU = {
    "mono": "DejaVu Sans Mono",
    "serif": "DejaVu Serif",
    "sans": "DejaVu Sans",
}
_CALIBRATION = {
    ("sans", False): 0.89,
    ("sans", True): 0.84,
    ("serif", False): 0.87,
    ("serif", True): 0.91,
    ("mono", False): 1.0,
    ("mono", True): 1.0,
}


def _font_kind(name: str | None) -> str:
    n = name or ""
    if _MONO_PAT.search(n):
        return "mono"
    if _SERIF_PAT.search(n):
        return "serif"
    return "sans"


@lru_cache(maxsize=64)
def _resolved_font(name: str | None, bold: bool, italic: bool):
    """(FontProperties at 100pt, width calibration).

    Prefers the *actual* declared font when matplotlib can find it installed
    (exact metrics); otherwise falls back to the bundled DejaVu family of the
    same kind with a measured width-calibration factor."""
    from matplotlib import font_manager as fm
    from matplotlib.font_manager import FontProperties

    kind = _font_kind(name)
    weight = "bold" if bold else "normal"
    style = "italic" if italic else "normal"
    if name:
        try:
            path = fm.findfont(
                FontProperties(family=name, weight=weight, style=style),
                fallback_to_default=False,
            )
            return FontProperties(fname=path, size=100), 1.0
        except Exception:
            pass
    return (
        FontProperties(family=_DEJAVU[kind], weight=weight, style=style, size=100),
        _CALIBRATION[(kind, bold)],
    )


@lru_cache(maxsize=4096)
def _width_at_100pt(text: str, name: str | None, bold: bool, italic: bool) -> float:
    """Width of ``text`` in points at 100pt (TextToPath works at dpi 72)."""
    if not text:
        return 0.0
    import warnings

    from matplotlib.textpath import TextToPath

    prop, calibration = _resolved_font(name, bold, italic)
    with warnings.catch_warnings():
        # a font missing a glyph (e.g. "→" in Georgia) measures with a
        # substitute width — close enough, and not worth log noise per call
        warnings.simplefilter("ignore")
        w, _h, _d = TextToPath().get_text_width_height_descent(text, prop, ismath=False)
    return float(w) * calibration


def text_width_in(
    text: str, font_name: str | None, size_pt: float, bold: bool, italic: bool
) -> float:
    """Rendered width of a run, in inches."""
    return (
        _width_at_100pt(text, font_name, bool(bold), bool(italic))
        * (size_pt / 100.0)
        / 72.0
    )


def _resolve_run(run, para) -> tuple[str | None, float, bool, bool]:
    """(font name, size pt, bold, italic) with paragraph-level fallback."""
    f, pf = run.font, para.font
    name = f.name or pf.name
    size = f.size or pf.size
    size_pt = size.pt if size is not None else 18.0
    bold = f.bold if f.bold is not None else bool(pf.bold)
    italic = f.italic if f.italic is not None else bool(pf.italic)
    return name, size_pt, bold, italic


def _para_metrics(para) -> tuple[float, float, int]:
    """(single-line width in, max font size pt, explicit ``<a:br/>`` count)."""
    width = 0.0
    max_sz = 0.0
    for run in para.runs:
        name, size_pt, bold, italic = _resolve_run(run, para)
        width += text_width_in(run.text, name, size_pt, bold, italic)
        max_sz = max(max_sz, size_pt)
    brs = len(
        para._p.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}br")
    )
    return width, max_sz or 18.0, brs


def _line_height_in(para, size_pt: float) -> float:
    factor = LINE_FACTOR
    ls = para.line_spacing
    if isinstance(ls, (int, float)) and not hasattr(ls, "pt"):
        factor *= float(ls)
    elif ls is not None and hasattr(ls, "pt"):
        return float(ls.pt) / 72.0
    return size_pt / 72.0 * factor


def _in(v) -> float:
    return (int(v) if v is not None else 0) / EMU_PER_IN


def _margins_in(tf) -> tuple[float, float, float, float]:
    def g(x, default):
        try:
            return _in(x)
        except Exception:
            return default

    return (
        g(tf.margin_left, 0.1),
        g(tf.margin_right, 0.1),
        g(tf.margin_top, 0.05),
        g(tf.margin_bottom, 0.05),
    )


def _vertical_overlap(a_top, a_bot, b_top, b_bot, eps: float = 0.03) -> bool:
    return (b_top < a_bot - eps) and (b_bot > a_top + eps)


def _room_in_direction(shape, others, slide_w_in: float, direction: str) -> float:
    """How far (inches) ``shape``'s edge can move in ``direction`` before hitting
    a neighbouring text shape (with a gap) or the slide edge."""
    left, top = _in(shape.left), _in(shape.top)
    right, bot = left + _in(shape.width), top + _in(shape.height)
    if direction == "right":
        limit = slide_w_in - EDGE_MARGIN_IN
        for o in others:
            ol = _in(o.left)
            if ol >= right - 0.02 and _vertical_overlap(
                top, bot, _in(o.top), _in(o.top) + _in(o.height)
            ):
                limit = min(limit, ol - NEIGHBOR_GAP_IN)
        return max(0.0, limit - right)
    limit = EDGE_MARGIN_IN
    for o in others:
        o_r = _in(o.left) + _in(o.width)
        if o_r <= left + 0.02 and _vertical_overlap(
            top, bot, _in(o.top), _in(o.top) + _in(o.height)
        ):
            limit = max(limit, o_r + NEIGHBOR_GAP_IN)
    return max(0.0, left - limit)


def _scale_runs(tf, factor: float) -> None:
    """Explicitly set every run's font size to ``size × factor`` (0.5pt steps)."""
    from pptx.util import Pt

    for para in tf.paragraphs:
        for run in para.runs:
            _name, size_pt, _b, _i = _resolve_run(run, para)
            new = max(1.0, math.floor(size_pt * factor * 2) / 2.0)
            run.font.size = Pt(new)


def fit_slide_text(slide, slide_w_in: float = 20.0) -> list[dict[str, Any]]:
    """Make every text shape on ``slide`` hold its text without renderer-side
    autofit: widen single-line labels (obstacle-aware), then shrink explicitly.

    Returns an audit list of the adjustments made (for tests/logging).
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu

    text_shapes = [
        sh
        for sh in slide.shapes
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
    ]
    audit: list[dict[str, Any]] = []

    for sh in text_shapes:
        if getattr(sh, "rotation", 0):
            continue
        tf = sh.text_frame
        paras = [p for p in tf.paragraphs if any(r.text.strip() for r in p.runs)]
        if not paras:
            continue
        ml, mr, mt, mb = _margins_in(tf)
        inner_w = _in(sh.width) - ml - mr
        inner_h = _in(sh.height) - mt - mb
        if inner_w <= 0.05 or inner_h <= 0.03:
            continue

        metrics = [_para_metrics(p) for p in paras]
        max_sz = max(m[1] for m in metrics)
        line_h = _line_height_in(paras[0], max_sz)
        single_line = (
            len(paras) == 1
            and metrics[0][2] == 0
            and inner_h < SINGLE_LINE_RATIO * line_h
        )

        if single_line:
            need = metrics[0][0] * SAFETY
            if need <= inner_w:
                continue
            others = [o for o in text_shapes if o is not sh]
            algn = paras[0].alignment
            grow = need - inner_w
            grew = 0.0
            if algn in (None, PP_ALIGN.LEFT, PP_ALIGN.JUSTIFY):
                room = _room_in_direction(sh, others, slide_w_in, "right")
                grew = min(grow, room)
                sh.width = Emu(int(sh.width) + int(grew * EMU_PER_IN))
            elif algn == PP_ALIGN.RIGHT:
                room = _room_in_direction(sh, others, slide_w_in, "left")
                grew = min(grow, room)
                sh.left = Emu(int(sh.left) - int(grew * EMU_PER_IN))
                sh.width = Emu(int(sh.width) + int(grew * EMU_PER_IN))
            elif algn == PP_ALIGN.CENTER:
                room_r = _room_in_direction(sh, others, slide_w_in, "right")
                room_l = _room_in_direction(sh, others, slide_w_in, "left")
                half = min(grow / 2, room_r, room_l)
                sh.left = Emu(int(sh.left) - int(half * EMU_PER_IN))
                sh.width = Emu(int(sh.width) + int(2 * half * EMU_PER_IN))
                grew = 2 * half
            inner_w_new = inner_w + grew
            entry = {
                "text": tf.text[:40],
                "mode": "single",
                "needed_in": round(need, 3),
                "grew_in": round(grew, 3),
                "scaled": None,
            }
            if need > inner_w_new + 1e-6:
                factor = max(MIN_SCALE, inner_w_new / need)
                _scale_runs(tf, factor)
                entry["scaled"] = round(factor, 3)
            audit.append(entry)
        else:
            total_h = 0.0
            for p, (w, sz, brs) in zip(paras, metrics):
                lines = max(1, math.ceil((w * SAFETY) / inner_w)) + brs
                total_h += lines * _line_height_in(p, sz)
            # designed-to-wrap boxes are cut tight; only act on clear overflow,
            # and leave estimated slack so near-fits are never touched twice
            if total_h <= inner_h * MULTI_TOLERANCE:
                continue
            factor = max(
                MIN_SCALE, min(1.0, math.sqrt(inner_h * MULTI_TARGET / total_h))
            )
            _scale_runs(tf, factor)
            audit.append(
                {
                    "text": tf.text[:40],
                    "mode": "multi",
                    "est_h_in": round(total_h, 3),
                    "inner_h_in": round(inner_h, 3),
                    "scaled": round(factor, 3),
                }
            )
    return audit


def fit_presentation_text(prs) -> list[dict[str, Any]]:
    """Run :func:`fit_slide_text` over every slide; returns the joined audit."""
    slide_w_in = _in(prs.slide_width) or 20.0
    audit: list[dict[str, Any]] = []
    for slide in prs.slides:
        audit.extend(fit_slide_text(slide, slide_w_in))
    return audit


__all__ = [
    "fit_slide_text",
    "fit_presentation_text",
    "text_width_in",
]
