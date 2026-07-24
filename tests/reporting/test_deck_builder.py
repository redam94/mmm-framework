"""python-pptx builder that fills the designed template from a fitted model.

Fast tests cover the bundled-template resolver and the low-level template helpers
(formatting-preserving text set, label-anchored card fill) on a tiny synthetic
deck. The slow test builds the real template from a fitted model and re-parses the
output: the headline KPI cards (with 80% ranges), the channel scorecard, the
ROI/decomposition image swaps, and one per-channel deep-dive slide per channel
(extras trimmed) each carrying its saturation/zone chart.
"""

from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from mmm_framework.reporting.deck import builder, template as T  # noqa: E402


def test_default_template_exists():
    p = builder.default_template_path()
    assert p.exists() and p.suffix == ".pptx"


def _mini_prs():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # a KPI card: label on top, value below, sub below that (same left)
    label = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.4))
    label.text_frame.text = "RETURN / $1"
    value = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(3), Inches(0.6))
    value.text_frame.text = "0.00"
    sub = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(3), Inches(0.4))
    sub.text_frame.text = "placeholder"
    return prs, slide


def test_set_text_preserves_single_run():
    _, slide = _mini_prs()
    sh = next(T.iter_text_shapes(slide))
    assert T.set_text(sh, "RETURN / $1 (edited)")
    assert sh.text_frame.text == "RETURN / $1 (edited)"


def test_fill_card_label_anchored():
    _, slide = _mini_prs()
    assert T.fill_card(slide, "RETURN / $1", "1.52", "80% 1.08–2.08")
    texts = [sh.text_frame.text for sh in T.iter_text_shapes(slide)]
    assert "RETURN / $1" in texts  # label kept
    assert "1.52" in texts and "80% 1.08–2.08" in texts
    assert "placeholder" not in texts  # sub overwritten


# ---------------------------------------------------------------------------
# text fitting: filled decks must not rely on renderer-side autofit
# (Google Slides never recomputes PPT's shrink-on-overflow)
# ---------------------------------------------------------------------------


def _single_line_overflows(sh) -> bool:
    """True when a single-line label's measured text is wider than its box —
    the condition that wraps/clips in Google Slides."""
    from mmm_framework.reporting.deck import textfit

    tf = sh.text_frame
    paras = [p for p in tf.paragraphs if any(r.text.strip() for r in p.runs)]
    if len(paras) != 1:
        return False
    w, max_sz, brs = textfit._para_metrics(paras[0])
    if brs:
        return False
    ml, mr, mt, mb = textfit._margins_in(tf)
    inner_h = textfit._in(sh.height) - mt - mb
    if inner_h >= textfit.SINGLE_LINE_RATIO * textfit._line_height_in(paras[0], max_sz):
        return False  # multi-line body box — wrapping is fine
    inner_w = textfit._in(sh.width) - ml - mr
    return w * textfit.SAFETY > inner_w + 0.01


def _assert_no_single_line_overflow(slides, context: str):
    bad = []
    for i, s in enumerate(slides, 1):
        for sh in s.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            if not sh.text_frame.text.strip() or getattr(sh, "rotation", 0):
                continue
            if _single_line_overflows(sh):
                bad.append((i, sh.text_frame.text[:40]))
    assert not bad, f"{context}: single-line text wider than its box: {bad}"


def test_fit_slide_text_grows_or_shrinks_long_labels():
    from pptx import Presentation
    from pptx.util import Inches

    from mmm_framework.reporting.deck import textfit

    prs = Presentation(str(builder.default_template_path()))
    # slide 13 (0-based 12): the deep-dive channel-name title, sized for "Video"
    slide = list(prs.slides)[12]
    title = next(
        sh for sh in T.iter_text_shapes(slide) if sh.text_frame.text.strip() == "Video"
    )
    T.set_text(title, "Programmatic Display")
    audit = textfit.fit_slide_text(slide, 20.0)
    assert any(a["text"].startswith("Programmatic") for a in audit)
    _assert_no_single_line_overflow([slide], "after fit")
    # the title either grew or the run shrank — and stayed inside the slide
    assert T._emu(title.left) + T._emu(title.width) <= int(Inches(20.0))


def test_fit_pass_fixes_untouched_template_labels():
    """Even the unfilled template has glyph-exact boxes ("Media" clips in
    Google Slides); the fit pass must leave no single-line overflow anywhere."""
    from pptx import Presentation

    from mmm_framework.reporting.deck import textfit

    prs = Presentation(str(builder.default_template_path()))
    textfit.fit_presentation_text(prs)
    _assert_no_single_line_overflow(prs.slides, "template after fit")


def test_template_has_no_run_together_boundaries():
    """The bundled template's label runs all carry their separating space (two
    were authored without it and patched); the builder normalizer is idempotent."""
    from pptx import Presentation

    prs = Presentation(str(builder.default_template_path()))
    for s in prs.slides:
        assert T.normalize_run_boundaries(s) == 0


def test_harvest_status_styles_reads_template_design():
    from pptx import Presentation

    prs = Presentation(str(builder.default_template_path()))
    styles = T.harvest_status_styles(prs)
    for status in T.READ_TO_ACTION:
        assert status in styles and styles[status]["text_rgb"] is not None
    # the four statuses use four distinct design colours
    colors = {str(styles[s]["text_rgb"]) for s in T.READ_TO_ACTION}
    assert len(colors) == 4
    # actions inherit pill decoration colours from the deep-dive pills
    assert styles["hold"]["bg_rgb"] is not None


def _row(ch, roi, lo, hi, action, spend=1000.0, hl=None):
    return {
        "channel": ch,
        "spend": spend,
        "contribution": roi * spend,
        "roi_mean": roi,
        "roi_lo": lo,
        "roi_hi": hi,
        "action": action,
        "half_life": hl,
    }


def test_headline_standfirst_never_funds_a_below_breakeven_channel():
    """An all-Reduce portfolio must not read "move budget toward <channel>"
    when that channel's own scorecard row says Reduce."""
    rows = [
        _row("X", 0.30, 0.10, 0.55, "Reduce"),
        _row("Y", 0.90, 0.60, 0.98, "Reduce"),
    ]
    text = builder._headline_standfirst({}, rows, be=1.0)
    assert "toward Y" not in text and "toward X" not in text
    assert "trimming X" in text
    # ... but a genuine Scale channel is recommended
    rows[1] = _row("Y", 1.60, 1.15, 2.1, "Scale")
    text = builder._headline_standfirst({}, rows, be=1.0)
    assert "toward Y" in text


def test_what_to_do_scale_does_not_overclaim_marginal_return():
    class _Z:
        current_mroi = 0.6  # below break-even despite a Scale-worthy average

    txt = builder._what_to_do_text(_row("A", 1.7, 1.2, 2.3, "Scale"), _Z(), be=1.0)
    assert "still clears break-even" not in txt
    assert "0.60" in txt
    _Z.current_mroi = 1.4
    txt = builder._what_to_do_text(_row("A", 1.7, 1.2, 2.3, "Scale"), _Z(), be=1.0)
    assert "still clears break-even" in txt


def test_fill_tests_slide_removes_unused_card_frames():
    from pptx import Presentation

    prs = Presentation(str(builder.default_template_path()))
    s22 = next(
        s for s in prs.slides if T.find_by_label(s, "RECOMMENDED TESTS") is not None
    )

    def _cards(slide):
        n = 0
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                continue
            try:
                prst = sh._element.spPr.prstGeom.attrib.get("prst")
            except Exception:
                continue
            # the triplet card frames are roundRects (blanked TEXT shapes are
            # plain rects — don't count those)
            if prst == "roundRect" and T._emu(sh.height) / 914400 < 2.0:
                n += 1
        return n

    before = _cards(s22)
    builder._fill_tests_slide(s22, [_row("Solo", 1.1, 0.4, 1.9, "Test")])
    after = _cards(s22)
    assert before == 3 and after == 1  # two unused triplets' card frames deleted
    texts = " ".join(sh.text_frame.text for sh in T.iter_text_shapes(s22))
    assert "Solo" in texts and "Geo holdout on Print" not in texts


def test_fill_recommend_cards_replaces_example_recommendations():
    from pptx import Presentation

    prs = Presentation(str(builder.default_template_path()))
    s2 = list(prs.slides)[1]
    rows = [
        _row("CTV", 1.6, 1.2, 2.1, "Scale"),
        _row("Audio", 0.5, 0.2, 0.8, "Reduce"),
        _row("OOH", 1.3, 0.6, 2.2, "Test"),
    ]
    builder._fill_recommend_cards(s2, rows, "$")
    texts = " ".join(sh.text_frame.text for sh in T.iter_text_shapes(s2))
    assert "Scale CTV" in texts and "Reduce Audio" in texts and "Test OOH" in texts
    assert "Scale Video" not in texts and "$18.9K of spend returning" not in texts


# ---------------------------------------------------------------------------
# slow: build the real template from a fitted model
# ---------------------------------------------------------------------------


@pytest.fixture(scope="module")
def fitted_model():
    from mmm_framework.config import ModelConfig
    from mmm_framework.model import BayesianMMM, TrendConfig
    from mmm_framework.model.trend_config import TrendType
    from mmm_framework.synth import dgp

    panel = dgp.build("clean", seed=0, n_weeks=104).panel()
    mmm = BayesianMMM(
        panel,
        ModelConfig(use_parametric_adstock=True),
        TrendConfig(type=TrendType.LINEAR),
    )
    mmm.fit(
        draws=250,
        tune=500,
        chains=2,
        target_accept=0.9,
        random_seed=3,
        progressbar=False,
    )
    return mmm


@pytest.mark.slow
def test_build_pptx_fills_template(fitted_model, tmp_path):
    from pptx import Presentation

    out = tmp_path / "deck.pptx"
    data = builder.build_pptx(
        fitted_model,
        out_path=out,
        client="Acme Corp",
        kpi_name="Sales",
        currency="$",
        break_even=1.0,
        hdi_prob=0.8,
    )
    assert isinstance(data, bytes) and len(data) > 10000 and out.exists()

    prs = Presentation(str(out))

    def texts(s):
        return [
            sh.text_frame.text
            for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]

    def has_label(s, lbl):
        return any(t.strip().lower() == lbl.lower() for t in texts(s))

    channels = set(fitted_model.channel_names)

    # the unused per-channel deep-dives are trimmed (template has 7; model has 4)
    deep = [
        s
        for s in prs.slides
        if has_label(s, "RETURN / $1") and has_label(s, "CARRYOVER HALF-LIFE")
    ]
    assert len(deep) == len(channels), (len(deep), len(channels))
    # each deep-dive names a distinct model channel and carries its zone chart
    named = []
    for s in deep:
        nm = next((t for t in texts(s) if t in channels), None)
        named.append(nm)
        npics = sum(1 for sh in s.shapes if "PICTURE" in str(sh.shape_type))
        assert npics >= 1, "saturation/zone chart should be inserted"
    assert set(named) == channels

    # headline KPI cards filled with 80% ranges
    head = next(s for s in prs.slides if has_label(s, "THE HEADLINE"))
    htexts = " ".join(texts(head))
    assert "80% range" in htexts
    # the S2 headline / WHAT WE RECOMMEND band never keep example-model claims
    assert "Marketing earns its keep" not in htexts
    assert "Scale Video" not in htexts and "$18.9K" not in htexts

    # the flighting-plan standfirst no longer names example channels
    s11 = next(s for s in prs.slides if has_label(s, "THE FLIGHTING PLAN"))
    s11_text = " ".join(texts(s11))
    assert "trims TV" not in s11_text and "Search's always-on" not in s11_text

    # scorecard names every model channel
    sc = next(s for s in prs.slides if has_label(s, "CHANNEL SCORECARD"))
    sc_texts = texts(sc)
    assert channels.issubset(set(sc_texts))

    # --- renderer-agnostic layout: nothing wraps or clips in Google Slides ---
    _assert_no_single_line_overflow(prs.slides, "built deck")
    for s in prs.slides:
        assert T.normalize_run_boundaries(s) == 0  # no run-together labels

    # --- READ pills coloured for the FILLED status, not the template row's ---
    from pptx import Presentation as _P

    tpl_styles = T.harvest_status_styles(_P(str(builder.default_template_path())))
    read_cells = [
        sh
        for sh in T.iter_text_shapes(sc)
        if T._norm(sh.text_frame.text) in T.READ_TO_ACTION
    ]
    assert read_cells, "scorecard should have filled READ cells"
    for sh in read_cells:
        status = T._norm(sh.text_frame.text)
        got = sh.text_frame.paragraphs[0].runs[0].font.color.rgb
        assert str(got) == str(
            tpl_styles[status]["text_rgb"]
        ), f"READ pill '{status}' keeps a stale colour"
        bg, dot = T.find_pill_parts(sc, sh)
        assert bg is not None
        assert str(T._fill_rgb(bg)) == str(tpl_styles[status]["bg_rgb"])

    # --- deep-dives carry channel-grounded prose, not the example model's ---
    stale_standfirsts = (
        "The portfolio's anchor",
        "The highest central return",
        "A large budget earning a thin return",
    )
    for s in deep:
        joined = " ".join(texts(s))
        assert not any(t in joined for t in stale_standfirsts), joined[:200]
        wtd = T.find_by_prefix(s, "What to do")
        assert wtd is not None
        body = wtd.text_frame.text
        assert "Increase weight and hold it continuously" not in body  # example text

    # --- S10 quadrants + S22 tests reference the model's channels ---
    s10 = next(s for s in prs.slides if has_label(s, "WHERE TO MOVE BUDGET"))
    s10_text = " ".join(texts(s10))
    assert not ("Print · Radio" in s10_text and "Print" not in channels)
    assert any(ch in s10_text for ch in channels)
    s22 = next(s for s in prs.slides if has_label(s, "RECOMMENDED TESTS"))
    s22_text = " ".join(texts(s22))
    assert "Geo holdout on Print" not in s22_text or "Print" in channels
    assert any(ch in s22_text for ch in channels)

    # --- decomposition waterfall image swapped in ---
    s4 = next(
        s for s in prs.slides if T.find_by_prefix(s, "The full revenue decomposition")
    )
    assert any("PICTURE" in str(sh.shape_type) for sh in s4.shapes)
