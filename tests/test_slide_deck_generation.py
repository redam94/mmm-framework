"""Agentic slide-deck generation: per-slide AI insights + whole-deck synthesis,
and the deterministic model-ops that outline and render the .pptx.

Fast tests drive the insight layer with a fake LLM (no network): correct keys,
call count, and graceful degradation. The slow test runs the real round-trip on a
fitted model — slide_deck_notes → insights (fake LLM) → render_slide_deck — and
confirms the .pptx is written with the insight text landed in the template.
"""

from __future__ import annotations

import os
import tempfile

import pytest

from mmm_framework.agents.deck_insights import generate_deck_insights


class _FakeResp:
    def __init__(self, content: str):
        self.content = content


class _FakeLLM:
    """Records calls and echoes a deterministic answer (no network)."""

    def __init__(self, fail: bool = False):
        self.calls = 0
        self.fail = fail

    def invoke(self, messages):
        self.calls += 1
        if self.fail:
            raise RuntimeError("LLM unavailable")
        # last message is the human prompt; tag the reply so we can find it
        return _FakeResp(
            f"INSIGHT-{self.calls}: scale where marginal ROI clears break-even."
        )


def _notes():
    return [
        {
            "key": "title",
            "kind": "title",
            "title": "Cover",
            "notes": "n",
            "is_summary": True,
            "metrics": {},
        },
        {
            "key": "headline",
            "kind": "executive_summary",
            "title": "Exec",
            "notes": "n",
            "is_summary": True,
            "metrics": {},
        },
        {
            "key": "channel:TV",
            "kind": "saturation",
            "title": "TV",
            "notes": "TV is in the saturation zone.",
            "is_summary": False,
            "metrics": {"zone": {"current_zone": "saturation"}},
        },
        {
            "key": "channel:Search",
            "kind": "saturation",
            "title": "Search",
            "notes": "Search under-invested.",
            "is_summary": False,
            "metrics": {"zone": {"current_zone": "breakthrough"}},
        },
        {
            "key": "reallocation",
            "kind": "optimization",
            "title": "Reallocate",
            "notes": "n",
            "is_summary": True,
            "metrics": {},
        },
    ]


# ---------------------------------------------------------------------------
# fast: the insight + synthesis layer
# ---------------------------------------------------------------------------


def test_generate_insights_keys_and_call_count():
    llm = _FakeLLM()
    ins = generate_deck_insights(_notes(), llm)
    # one insight per channel deep-dive + one synthesized headline
    assert set(ins) == {"channel:TV", "channel:Search", "headline"}
    assert llm.calls == 3  # 2 channels + 1 synthesis
    assert all(ins.values())


def test_generate_insights_degrades_on_llm_failure():
    ins = generate_deck_insights(_notes(), _FakeLLM(fail=True))
    assert ins == {}  # no crash, no narrative


def test_generate_insights_skips_summary_slides_for_per_slide_pass():
    # only saturation (channel) slides get a per-slide insight; summaries are
    # handled by the synthesis call (the "headline" key)
    llm = _FakeLLM()
    ins = generate_deck_insights(_notes(), llm)
    assert "title" not in ins and "reallocation" not in ins


# ---------------------------------------------------------------------------
# slow: model-ops round-trip on a fitted model
# ---------------------------------------------------------------------------


@pytest.mark.slow
def test_ops_roundtrip_renders_pptx_with_insights():
    from pptx import Presentation

    from mmm_framework.agents import model_ops
    from mmm_framework.agents.runtime import set_current_thread
    from mmm_framework.config import ModelConfig
    from mmm_framework.model import BayesianMMM, TrendConfig
    from mmm_framework.model.trend_config import TrendType
    from mmm_framework.synth import dgp

    os.environ["MMM_AGENT_WORKSPACE"] = tempfile.mkdtemp()
    set_current_thread("__deck_roundtrip__")

    panel = dgp.build("clean", seed=0, n_weeks=104).panel()
    mmm = BayesianMMM(
        panel,
        ModelConfig(use_parametric_adstock=True),
        TrendConfig(type=TrendType.LINEAR),
    )
    mmm.fit(
        draws=200,
        tune=400,
        chains=2,
        target_accept=0.9,
        random_seed=3,
        progressbar=False,
    )

    # 1) outline
    r1 = model_ops.slide_deck_notes(mmm, None, client="Acme", kpi_name="Sales")
    assert r1["error"] is None
    notes = r1["dashboard"]["slide_deck_notes"]
    assert any(n["key"].startswith("channel:") for n in notes)
    assert any(n["key"] == "headline" for n in notes)

    # 2) insights (fake LLM)
    llm = _FakeLLM()
    insights = generate_deck_insights(notes, llm)
    assert "headline" in insights

    # 3) render with insights
    r2 = model_ops.render_slide_deck(
        mmm, None, insights=insights, client="Acme", kpi_name="Sales"
    )
    assert r2["error"] is None
    path = r2["dashboard"]["slide_deck"]["path"]
    assert os.path.getsize(path) > 10000

    # the synthesized headline landed on the template's headline slide
    prs = Presentation(path)

    def texts(s):
        return [
            sh.text_frame.text
            for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]

    all_text = " ".join(t for s in prs.slides for t in texts(s))
    assert "INSIGHT-" in all_text  # at least one AI insight string is in the deck
